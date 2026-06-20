"""Power BI routes — chat page + API endpoints.

Falls back to demo mode with realistic mock data when the Power BI API
credentials are not configured.
"""

import json
import logging
import tempfile

from flask import Blueprint, jsonify, render_template, request, current_app, send_file

from src.utils.logging_utils import log_web_activity, get_client_ip
from src.components.web import powerbi_chat_handler as pbi_chat
from web.extensions import limiter

logger = logging.getLogger(__name__)
powerbi_bp = Blueprint("powerbi", __name__)

# Lazy-init Power BI client and LLM
_pbi_client = None
_pbi_llm = None
_demo_mode = None  # True/False, set on first _get_pbi_client call


def _is_demo_mode() -> bool:
    global _demo_mode
    if _demo_mode is None:
        _get_pbi_client()  # triggers detection
    return _demo_mode


def _get_pbi_client():
    global _pbi_client, _demo_mode
    if _pbi_client is None:
        try:
            from services.powerbi import PowerBIClient
            _pbi_client = PowerBIClient()
            _demo_mode = False
        except Exception:
            from services.powerbi_demo import DemoClient
            _pbi_client = DemoClient()
            _demo_mode = True
            logger.info("Power BI API not configured — using demo mode")
    return _pbi_client


def _get_pbi_llm():
    global _pbi_llm
    if _pbi_llm is None:
        from my_bot.utils.llm_factory import create_llm
        _pbi_llm = create_llm(
            max_tokens=2048, timeout=300,
            extra_body={"chat_template_kwargs": {"enable_thinking": False}},
        )
    return _pbi_llm


_prune_llm = None


def _get_prune_llm():
    """Lightweight LLM for schema pruning — Llama-3.2-1B on separate process."""
    global _prune_llm
    if _prune_llm is None:
        import os
        import httpx
        # M1 Router: Qwen3-8B-4bit (port 8016) — on M1 Mac's GPU, zero M3 contention.
        base_url = os.environ.get("POWERBI_PRUNE_LLM_URL") or "http://localhost:8016/v1"
        # Discover model name from the server
        model_name = "default"
        try:
            resp = httpx.get(f"{base_url}/models", timeout=3)
            models = resp.json().get("data", [])
            if models:
                model_name = models[0]["id"]
        except Exception:
            pass
        from langchain_openai import ChatOpenAI
        _prune_llm = ChatOpenAI(
            base_url=base_url, model=model_name, api_key="not-needed",
            temperature=0, max_tokens=200, timeout=10,
            extra_body={"enable_thinking": False},
        )
        logger.info("Schema prune LLM: %s model=%s", base_url, model_name)
    return _prune_llm


# ── Page route ──

@powerbi_bp.route("/powerbi")
@log_web_activity
def powerbi_page():
    """Render the Power BI chat page."""
    datasets = []
    try:
        client = _get_pbi_client()
        datasets = client.list_datasets()
    except Exception as exc:
        logger.warning("Power BI datasets unavailable: %s", exc)
    return render_template("powerbi_chat.html", datasets=datasets, demo_mode=_is_demo_mode())


# ── API endpoints ──

@powerbi_bp.route("/api/powerbi/datasets")
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_datasets():
    """List available Power BI datasets."""
    try:
        client = _get_pbi_client()
        datasets = client.list_datasets()
        return jsonify({"success": True, "datasets": datasets})
    except Exception as exc:
        logger.error("Power BI datasets error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@powerbi_bp.route("/api/powerbi/refresh/<dataset_id>")
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_refresh(dataset_id):
    """Get last refresh time for a dataset."""
    try:
        client = _get_pbi_client()
        refresh = client.get_last_refresh(dataset_id)
        if refresh:
            return jsonify({"success": True, "refresh": refresh})
        return jsonify({"success": True, "refresh": None})
    except Exception as exc:
        logger.warning("Power BI refresh check error: %s", exc)
        return jsonify({"success": True, "refresh": None})


@powerbi_bp.route("/api/powerbi/charts/<dataset_id>")
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_charts(dataset_id):
    """Return chart data, KPIs, and suggestion chips for a dataset."""
    try:
        if _is_demo_mode():
            from services.powerbi_demo import get_chart_data
            data = get_chart_data(dataset_id)
            if not data:
                return jsonify({"success": False, "error": "Unknown dataset"}), 404
            return jsonify({"success": True, **data})

        # Live mode: auto-generate charts from schema
        from services.powerbi_chart_builder import build_charts_and_chips
        client = _get_pbi_client()
        dataset_name = request.args.get("name", "")
        schema_result = client.execute_dax(dataset_id, "EVALUATE COLUMNSTATISTICS()")
        if schema_result.get("error"):
            return jsonify({"success": False, "error": schema_result["error"]}), 400
        data = build_charts_and_chips(client, dataset_id, schema_result["rows"], dataset_name=dataset_name)
        return jsonify({"success": True, **data})
    except Exception as exc:
        logger.error("Power BI charts error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


_schema_cache: dict[str, tuple[str, float]] = {}  # dataset_id -> (full_schema, timestamp)
_SCHEMA_TTL = 3600  # 1 hour


@powerbi_bp.route("/api/powerbi/schema/<dataset_id>")
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_schema(dataset_id):
    """Fetch and cache schema for a dataset."""
    try:
        # Demo mode: return pre-built schema directly
        if _is_demo_mode():
            from services.powerbi_demo import DEMO_SCHEMAS
            schema_text = DEMO_SCHEMAS.get(dataset_id, "Schema not available.")
            pbi_chat.set_dataset_schema(dataset_id, schema_text)
            return jsonify({"success": True, "schema": schema_text})

        # Return server-cached schema if fresh
        import time as _time
        cached = _schema_cache.get(dataset_id)
        if cached and (_time.time() - cached[1]) < _SCHEMA_TTL:
            # Ensure LLM schema is set (may have been lost on restart)
            if not pbi_chat.get_dataset_schema(dataset_id).startswith("No schema"):
                return jsonify({"success": True, "schema": cached[0]})

        client = _get_pbi_client()
        # COLUMNSTATISTICS works on all dataset types
        result = client.execute_dax(dataset_id, "EVALUATE COLUMNSTATISTICS()")

        if result.get("error"):
            return jsonify({"success": False, "error": result["error"]}), 400

        # Build schema text from COLUMNSTATISTICS rows.
        # Two versions: full (for UI display) and compact (for LLM prompt).
        # LLM version: top 5 tables by column count, no ranges, keeps prompt small.
        MAX_LLM_TABLES = 5

        tables: dict[str, list[dict]] = {}  # tbl -> [{col, card, min, max}]
        for row in result.get("rows", []):
            tbl = row.get("[Table Name]") or row.get("Table") or ""
            col = row.get("[Column Name]") or row.get("Column") or ""
            if not tbl or not col:
                continue
            if "DateTableTemplate" in tbl or "LocalDateTable" in tbl:
                continue
            if "RowNumber-" in col:
                continue
            tables.setdefault(tbl, []).append({
                "col": col,
                "card": row.get("[Cardinality]", ""),
                "min": row.get("[Min]", ""),
                "max": row.get("[Max]", ""),
            })

        # Full schema for UI display
        full_lines = []
        for tbl, cols in sorted(tables.items()):
            full_lines.append(f"Table: {tbl}")
            for c in cols:
                hints = []
                if c["card"]:
                    hints.append(f"{c['card']:,} distinct" if isinstance(c["card"], int) else str(c["card"]))
                if c["min"] not in (None, "") and c["max"] not in (None, ""):
                    hints.append(f"range: {c['min']} .. {c['max']}")
                full_lines.append(f"  - {c['col']}" + (f" ({', '.join(hints)})" if hints else ""))
            full_lines.append("")
        full_schema = "\n".join(full_lines) or "Schema could not be determined."

        # Compact schema for LLM — top tables, max 30 cols each, ~150 lines target
        MAX_LLM_COLS_PER_TABLE = 30
        top_tables = sorted(tables.keys(), key=lambda t: len(tables[t]), reverse=True)[:MAX_LLM_TABLES]
        compact_lines = []
        for tbl in top_tables:
            cols = tables[tbl]
            compact_lines.append(f"Table: {tbl} ({len(cols)} columns)")
            shown = cols[:MAX_LLM_COLS_PER_TABLE]
            for c in shown:
                card_hint = f" ({c['card']:,} distinct)" if isinstance(c["card"], int) and c["card"] > 1 else ""
                compact_lines.append(f"  - {c['col']}{card_hint}")
            if len(cols) > MAX_LLM_COLS_PER_TABLE:
                compact_lines.append(f"  ... and {len(cols) - MAX_LLM_COLS_PER_TABLE} more columns")
            compact_lines.append("")
        if len(tables) > MAX_LLM_TABLES:
            skipped = sorted(set(tables.keys()) - set(top_tables))
            compact_lines.append(f"({len(skipped)} smaller tables not shown: {', '.join(skipped[:10])}{'...' if len(skipped) > 10 else ''})")
        compact_schema = "\n".join(compact_lines)

        # Cache compact version for the chat handler (keeps LLM prompt small)
        pbi_chat.set_dataset_schema(dataset_id, compact_schema)
        _schema_cache[dataset_id] = (full_schema, _time.time())

        return jsonify({"success": True, "schema": full_schema})
    except Exception as exc:
        logger.error("Power BI schema error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@powerbi_bp.route("/api/powerbi/route", methods=["POST"])
@limiter.limit("20 per minute")
@log_web_activity
def api_powerbi_route():
    """Cross-dataset router: NL question -> best dataset to answer it.

    Returns {success, dataset_id, dataset_name, confidence, reason,
    alternatives:[{id,name}], method}. Used by the Explorer's Auto mode so the
    user doesn't have to pick a dataset before asking.
    """
    try:
        data = request.get_json(silent=True) or {}
        question = (data.get("message") or "").strip()
        if not question:
            return jsonify({"success": False, "error": "Message is required"}), 400
        if len(question) > 2000:
            return jsonify({"success": False, "error": "Message too long"}), 400

        from services.powerbi_router import build_catalog, route_question

        if _is_demo_mode():
            # Demo mode has no live catalog — route over the demo dataset list.
            client = _get_pbi_client()
            catalog = build_catalog(client)
            result = route_question(question, catalog, None, history=data.get("history"))
            return jsonify({"success": True, **result})

        client = _get_pbi_client()
        catalog = build_catalog(client)
        # Prefer the lightweight router LLM (fast classification); fall back to
        # the main chat LLM, then to deterministic keyword routing inside route_question.
        try:
            router_llm = _get_prune_llm()
        except Exception:
            router_llm = None
        if router_llm is None:
            router_llm = _get_pbi_llm()
        result = route_question(question, catalog, router_llm, history=data.get("history"))
        return jsonify({"success": True, **result})
    except Exception as exc:
        logger.error("Power BI route error: %s", exc, exc_info=True)
        return jsonify({"success": False, "error": str(exc)}), 500


@powerbi_bp.route("/api/powerbi/chat/stream", methods=["POST"])
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_chat_stream():
    """Streaming chat: NL -> DAX -> execute -> explain."""
    try:
        data = request.get_json()
        user_message = (data.get("message") or "").strip()
        dataset_id = (data.get("dataset_id") or "").strip()
        session_id = (data.get("session_id") or "").strip()
        dataset_name = (data.get("dataset_name") or "").strip()
        client_history = data.get("history")  # list of {role, text} from client localStorage

        if not user_message:
            return jsonify({"success": False, "error": "Message is required"}), 400
        if len(user_message) > 2000:
            return jsonify({"success": False, "error": "Message too long (max 2000 chars)"}), 400
        if not dataset_id:
            return jsonify({"success": False, "error": "Dataset is required"}), 400
        if not session_id:
            return jsonify({"success": False, "error": "Session ID is required"}), 400

        llm = _get_pbi_llm()
        client = _get_pbi_client()
        client_ip = get_client_ip()
        prune_llm = _get_prune_llm()

        def generate():
            try:
                for payload in pbi_chat.handle_chat_stream(
                    user_message, dataset_id, session_id, llm, client,
                    client_ip=client_ip, dataset_name=dataset_name,
                    prune_llm=prune_llm, history=client_history,
                ):
                    if payload.get("keepalive"):
                        yield ": keepalive\n\n"
                    else:
                        yield f"data: {json.dumps(payload)}\n\n"
            except Exception as err:
                logger.error("Power BI chat stream error: %s", err, exc_info=True)
                err_str = str(err)
                if "incomplete chunked read" in err_str or "RemoteProtocolError" in err_str:
                    msg = "LLM connection dropped mid-response — the inference server may be overloaded. Please try again."
                else:
                    msg = err_str
                yield f"data: {json.dumps({'error': msg})}\n\n"

        return current_app.response_class(
            generate(),
            mimetype="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )
    except Exception as exc:
        logger.error("Power BI chat error: %s", exc, exc_info=True)
        return jsonify({"success": False, "error": str(exc)}), 500


@powerbi_bp.route("/api/powerbi/chips/<dataset_id>")
@limiter.limit("5 per minute")
@log_web_activity
def api_powerbi_chips(dataset_id):
    """Generate LLM-powered suggestion chips for a dataset (cached)."""
    try:
        dataset_name = request.args.get("name", "")

        # Check curated chips first (instant)
        from services.powerbi_chart_builder import CURATED_CHIPS
        if dataset_name:
            name_lower = dataset_name.lower()
            for key, curated in CURATED_CHIPS.items():
                if key in name_lower:
                    return jsonify({"success": True, "chips": curated})

        # Generate via LLM (cached after first call)
        llm = _get_pbi_llm()
        chips = pbi_chat.generate_llm_chips(llm, dataset_id, dataset_name)
        if chips:
            return jsonify({"success": True, "chips": chips})
        return jsonify({"success": True, "chips": []})
    except Exception as exc:
        logger.warning("Power BI chips error: %s", exc)
        return jsonify({"success": True, "chips": []})


@powerbi_bp.route("/api/powerbi/chat/clear", methods=["POST"])
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_chat_clear():
    """Clear Power BI chat session history."""
    data = request.get_json(silent=True) or {}
    session_id = (data.get("session_id") or "").strip()
    if not session_id:
        return jsonify({"success": False, "error": "Session ID is required"}), 400
    pbi_chat.clear_history(session_id)
    return jsonify({"success": True})


@powerbi_bp.route("/api/powerbi/export/xlsx", methods=["POST"])
@limiter.limit("10 per minute")
@log_web_activity
def api_powerbi_export_xlsx():
    """Export query results as a professionally formatted Excel file."""
    try:
        import pandas as pd
        from src.utils.excel_formatting import apply_professional_formatting

        data = request.get_json()
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        dataset_name = data.get("dataset_name", "PowerBI")

        if not headers or not rows:
            return jsonify({"success": False, "error": "No data to export"}), 400

        df = pd.DataFrame(rows, columns=headers)

        # Auto-detect column widths based on content
        col_widths = {}
        for col in headers:
            max_len = max(len(str(col)), df[col].astype(str).str.len().max())
            col_widths[col.lower()] = min(max(max_len + 4, 12), 60)

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        df.to_excel(tmp_path, index=False, sheet_name=dataset_name[:31])
        apply_professional_formatting(tmp_path, column_widths=col_widths)

        safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in dataset_name)
        filename = f"PowerBI - {safe_name}.xlsx"

        return send_file(
            tmp_path,
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            as_attachment=True,
            download_name=filename,
        )
    except Exception as exc:
        logger.error("Power BI Excel export error: %s", exc, exc_info=True)
        return jsonify({"success": False, "error": str(exc)}), 500


@powerbi_bp.route("/api/powerbi/export/pptx", methods=["POST"])
@limiter.limit("5 per minute")
@log_web_activity
def api_powerbi_export_pptx():
    """Export dashboard KPIs and chart summaries as a formatted PPTX."""
    try:
        from pathlib import Path
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        data = request.get_json()
        dataset_name = data.get("dataset_name", "Power BI Dataset")
        kpis = data.get("kpis", [])
        charts = data.get("charts", [])

        # Colors
        NAVY = RGBColor(0x1A, 0x23, 0x7E)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        BLUE = RGBColor(0x00, 0x46, 0xAD)
        GOLD = RGBColor(0xF6, 0xBE, 0x00)
        GRAY = RGBColor(0x64, 0x74, 0x8B)
        OFFWHITE = RGBColor(0xF8, 0xF9, 0xFF)

        # Try template
        template_path = Path("data/transient/PPTX Templates/Cybersecurity Slide Template_January2026.pptx")
        if template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # Find blank layout
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                blank_layout = layout
                break

        # Remove template slides
        while len(prs.slides._sldIdLst):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE

        # Title bar
        title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = NAVY
        title_shape.line.fill.background()

        title_tf = title_shape.text_frame
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.text = f"  \U0001F4CA  {dataset_name.replace('_', ' ')}"
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Gold accent bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), prs.slide_width, Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD
        accent.line.fill.background()

        # KPI cards
        if kpis:
            card_width = min(Inches(2.8), (prs.slide_width - Inches(1.6)) // len(kpis))
            start_x = Inches(0.6)
            for i, kpi in enumerate(kpis):
                x = start_x + i * (card_width + Inches(0.2))
                y = Inches(1.5)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, Inches(1.2))
                card.fill.solid()
                card.fill.fore_color.rgb = OFFWHITE
                card.line.color.rgb = BLUE
                card.line.width = Pt(1.5)

                tf = card.text_frame
                tf.word_wrap = True
                pv = tf.paragraphs[0]
                pv.text = kpi.get("value", "")
                pv.font.size = Pt(26)
                pv.font.bold = True
                pv.font.color.rgb = NAVY
                pv.alignment = PP_ALIGN.CENTER

                pl = tf.add_paragraph()
                pl.text = kpi.get("label", "")
                pl.font.size = Pt(10)
                pl.font.color.rgb = GRAY
                pl.font.bold = True
                pl.alignment = PP_ALIGN.CENTER

        # Chart summaries
        if charts:
            y_pos = Inches(3.0)
            for i, chart in enumerate(charts[:8]):
                col = i % 4
                row = i // 4
                x = Inches(0.6) + col * Inches(3.1)
                y = y_pos + row * Inches(2.0)
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(1.8))
                box.fill.solid()
                box.fill.fore_color.rgb = OFFWHITE
                box.line.color.rgb = RGBColor(0xC5, 0xCA, 0xE9)
                box.line.width = Pt(1)

                tf = box.text_frame
                tf.word_wrap = True
                pt = tf.paragraphs[0]
                pt.text = chart.get("title", "Chart")
                pt.font.size = Pt(12)
                pt.font.bold = True
                pt.font.color.rgb = NAVY
                pt.alignment = PP_ALIGN.LEFT

                if chart.get("insight"):
                    pi = tf.add_paragraph()
                    pi.text = chart["insight"]
                    pi.font.size = Pt(9)
                    pi.font.italic = True
                    pi.font.color.rgb = RGBColor(0x6A, 0x1B, 0x9A)

                # Show top labels/values
                labels = chart.get("labels", [])[:5]
                ds = chart.get("datasets", [{}])[0]
                vals = ds.get("data", [])[:5]
                for j, lbl in enumerate(labels):
                    pd_ = tf.add_paragraph()
                    val = vals[j] if j < len(vals) else ""
                    pd_.text = f"  {lbl}: {val}"
                    pd_.font.size = Pt(8)
                    pd_.font.color.rgb = GRAY

        # Watermark
        wm = slide.shapes.add_textbox(prs.slide_width - Inches(3.2), prs.slide_height - Inches(0.4), Inches(3), Inches(0.3))
        wm_tf = wm.text_frame
        wm_p = wm_tf.paragraphs[0]
        wm_p.text = "Power BI Explorer  \u2022  - Vinay Vobbilichetty"
        wm_p.font.size = Pt(7)
        wm_p.font.italic = True
        wm_p.font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)
        wm_p.alignment = PP_ALIGN.RIGHT

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        prs.save(tmp_path)

        safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in dataset_name)
        return send_file(
            tmp_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"PowerBI - {safe_name} Dashboard.pptx",
        )
    except Exception as exc:
        logger.error("Power BI PPTX export error: %s", exc, exc_info=True)
        return jsonify({"success": False, "error": str(exc)}), 500
